"""
Eva Live Document Processing Module

This module handles the ingestion and processing of various document formats
including PowerPoint presentations, PDFs, and text documents. It extracts
content, creates semantic chunks, and prepares documents for embedding.
"""

import asyncio
import logging
import time
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
from pathlib import Path
import hashlib
import re
import io

# Document processing libraries
from pptx import Presentation
import PyPDF2
import pytesseract
from PIL import Image

# Text processing
import tiktoken

from ..shared.config import get_config
from ..shared.models import PerformanceMetric

@dataclass
class DocumentChunk:
    """A semantic chunk of document content"""
    id: str
    content: str
    metadata: Dict[str, Any]
    chunk_index: int
    total_chunks: int
    token_count: int
    source_file: str
    source_type: str
    
@dataclass
class ProcessedDocument:
    """Processed document with chunks and metadata"""
    document_id: str
    title: str
    content: str
    chunks: List[DocumentChunk]
    metadata: Dict[str, Any]
    processing_time_ms: int
    total_tokens: int

class DocumentProcessor:
    """Document processing engine for Eva Live"""
    
    def __init__(self):
        self.config = get_config()
        self.logger = logging.getLogger(__name__)
        
        # Configuration
        self.max_chunk_size = self.config.get('content.processing.chunk_size', 1000)
        self.chunk_overlap = self.config.get('content.processing.overlap', 200)
        self.max_file_size = self.config.get('content.processing.max_file_size', '500MB')
        
        # Initialize tokenizer for chunk sizing
        self.tokenizer = tiktoken.get_encoding("cl100k_base")
        
        # Supported formats
        self.supported_formats = {
            '.pptx': self._process_powerpoint,
            '.ppt': self._process_powerpoint,
            '.pdf': self._process_pdf,
            '.txt': self._process_text,
            '.md': self._process_text,
            '.docx': self._process_docx
        }
        
        # Performance tracking
        self.metrics: List[PerformanceMetric] = []
    
    async def process_document(self, file_path: str, title: Optional[str] = None) -> ProcessedDocument:
        """Process a document and return chunks with metadata"""
        start_time = time.time()
        
        try:
            # Validate file
            path = Path(file_path)
            if not path.exists():
                raise FileNotFoundError(f"Document not found: {file_path}")
            
            file_extension = path.suffix.lower()
            if file_extension not in self.supported_formats:
                raise ValueError(f"Unsupported file format: {file_extension}")
            
            # Check file size
            file_size = path.stat().st_size
            max_size_bytes = self._parse_size(self.max_file_size)
            if file_size > max_size_bytes:
                raise ValueError(f"File too large: {file_size} bytes (max: {max_size_bytes})")
            
            # Generate document ID
            document_id = self._generate_document_id(file_path)
            
            # Extract content
            processor = self.supported_formats[file_extension]
            content, metadata = await processor(file_path)
            
            # Create semantic chunks
            chunks = await self._create_chunks(content, document_id, file_path, file_extension)
            
            processing_time = int((time.time() - start_time) * 1000)
            
            # Calculate total tokens
            total_tokens = sum(chunk.token_count for chunk in chunks)
            
            # Create processed document
            processed_doc = ProcessedDocument(
                document_id=document_id,
                title=title or path.stem,
                content=content,
                chunks=chunks,
                metadata={
                    **metadata,
                    'file_path': file_path,
                    'file_size': file_size,
                    'file_extension': file_extension,
                    'processed_at': time.time(),
                    'chunk_count': len(chunks)
                },
                processing_time_ms=processing_time,
                total_tokens=total_tokens
            )
            
            # Record metrics
            await self._record_metric("processing_time_ms", processing_time, "document_processor")
            await self._record_metric("chunk_count", len(chunks), "document_processor")
            await self._record_metric("total_tokens", total_tokens, "document_processor")
            
            self.logger.info(f"Processed document {document_id}: {len(chunks)} chunks, {total_tokens} tokens")
            
            return processed_doc
            
        except Exception as e:
            self.logger.error(f"Error processing document {file_path}: {e}")
            raise
    
    async def _process_powerpoint(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            
            content_parts = []
            slide_data = []
            
            for i, slide in enumerate(prs.slides):
                slide_content = []
                slide_title = ""
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        slide_content.append(text)
                        
                        # Try to identify slide title (usually first or largest text)
                        if not slide_title and len(text) < 100:
                            slide_title = text
                
                # Combine slide content
                slide_text = "\n".join(slide_content)
                if slide_text:
                    content_parts.append(f"=== Slide {i+1}: {slide_title or 'Untitled'} ===\n{slide_text}")
                    
                    slide_data.append({
                        'slide_number': i + 1,
                        'title': slide_title,
                        'content': slide_text,
                        'shape_count': len(slide.shapes)
                    })
            
            content = "\n\n".join(content_parts)
            
            metadata = {
                'slide_count': len(prs.slides),
                'slides': slide_data,
                'format': 'powerpoint'
            }
            
            return content, metadata
            
        except Exception as e:
            self.logger.error(f"Error processing PowerPoint {file_path}: {e}")
            raise
    
    async def _process_pdf(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process PDF document"""
        try:
            content_parts = []
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        text = page.extract_text()
                        if text.strip():
                            content_parts.append(f"=== Page {page_num + 1} ===\n{text}")
                    except Exception as e:
                        self.logger.warning(f"Error extracting text from page {page_num + 1}: {e}")
                        continue
            
            content = "\n\n".join(content_parts)
            
            metadata = {
                'page_count': len(pdf_reader.pages),
                'format': 'pdf'
            }
            
            return content, metadata
            
        except Exception as e:
            self.logger.error(f"Error processing PDF {file_path}: {e}")
            raise
    
    async def _process_text(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process plain text or markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            metadata = {
                'format': 'text',
                'encoding': 'utf-8'
            }
            
            return content, metadata
            
        except Exception as e:
            self.logger.error(f"Error processing text file {file_path}: {e}")
            raise
    
    async def _process_docx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process Word document (basic implementation)"""
        try:
            # Note: This would require python-docx package
            # For now, return empty content with warning
            self.logger.warning(f"DOCX processing not fully implemented for {file_path}")
            
            metadata = {
                'format': 'docx',
                'status': 'not_implemented'
            }
            
            return "", metadata
            
        except Exception as e:
            self.logger.error(f"Error processing DOCX {file_path}: {e}")
            raise
    
    async def _create_chunks(self, content: str, document_id: str, file_path: str, file_extension: str) -> List[DocumentChunk]:
        """Create semantic chunks from document content"""
        if not content.strip():
            return []
        
        chunks = []
        
        # Split content into sentences
        sentences = self._split_into_sentences(content)
        
        current_chunk = ""
        current_tokens = 0
        chunk_index = 0
        
        for sentence in sentences:
            sentence_tokens = len(self.tokenizer.encode(sentence))
            
            # Check if adding this sentence would exceed chunk size
            if current_tokens + sentence_tokens > self.max_chunk_size and current_chunk:
                # Create chunk
                chunk = await self._create_chunk(
                    current_chunk.strip(),
                    chunk_index,
                    document_id,
                    file_path,
                    file_extension
                )
                chunks.append(chunk)
                
                # Start new chunk with overlap
                overlap_text = self._get_overlap_text(current_chunk, self.chunk_overlap)
                current_chunk = overlap_text + sentence
                current_tokens = len(self.tokenizer.encode(current_chunk))
                chunk_index += 1
            else:
                current_chunk += " " + sentence if current_chunk else sentence
                current_tokens += sentence_tokens
        
        # Add final chunk
        if current_chunk.strip():
            chunk = await self._create_chunk(
                current_chunk.strip(),
                chunk_index,
                document_id,
                file_path,
                file_extension
            )
            chunks.append(chunk)
        
        # Update total chunks count
        for chunk in chunks:
            chunk.total_chunks = len(chunks)
        
        return chunks
    
    async def _create_chunk(self, content: str, chunk_index: int, document_id: str, file_path: str, file_extension: str) -> DocumentChunk:
        """Create a single document chunk"""
        chunk_id = f"{document_id}_chunk_{chunk_index}"
        token_count = len(self.tokenizer.encode(content))
        
        return DocumentChunk(
            id=chunk_id,
            content=content,
            metadata={
                'document_id': document_id,
                'source_file': file_path,
                'source_type': file_extension[1:],  # Remove dot
                'created_at': time.time()
            },
            chunk_index=chunk_index,
            total_chunks=0,  # Will be updated later
            token_count=token_count,
            source_file=file_path,
            source_type=file_extension[1:]
        )
    
    def _split_into_sentences(self, text: str) -> List[str]:
        """Split text into sentences using regex"""
        # Simple sentence splitting - could be improved with spaCy or NLTK
        sentences = re.split(r'(?<=[.!?])\s+', text)
        
        # Filter out very short "sentences"
        sentences = [s.strip() for s in sentences if len(s.strip()) > 10]
        
        return sentences
    
    def _get_overlap_text(self, text: str, overlap_chars: int) -> str:
        """Get overlap text from the end of current chunk"""
        if len(text) <= overlap_chars:
            return text
        
        # Try to find a sentence boundary for clean overlap
        overlap_text = text[-overlap_chars:]
        
        # Find the first sentence boundary
        sentences = self._split_into_sentences(overlap_text)
        if sentences:
            return sentences[0]
        
        return overlap_text
    
    def _generate_document_id(self, file_path: str) -> str:
        """Generate unique document ID based on file path and modification time"""
        path = Path(file_path)
        mtime = path.stat().st_mtime
        content = f"{file_path}_{mtime}"
        return hashlib.md5(content.encode()).hexdigest()
    
    def _parse_size(self, size_str: str) -> int:
        """Parse size string like '500MB' to bytes"""
        size_str = size_str.upper().strip()
        
        multipliers = {
            'B': 1,
            'KB': 1024,
            'MB': 1024 ** 2,
            'GB': 1024 ** 3
        }
        
        for suffix, multiplier in multipliers.items():
            if size_str.endswith(suffix):
                number = size_str[:-len(suffix)]
                return int(float(number) * multiplier)
        
        # Default to bytes
        return int(size_str)
    
    async def _record_metric(self, metric_type: str, value: float, component: str) -> None:
        """Record performance metric"""
        metric = PerformanceMetric(
            metric_type=metric_type,
            metric_value=value,
            component=component
        )
        self.metrics.append(metric)
    
    def get_metrics(self) -> List[PerformanceMetric]:
        """Get recorded metrics"""
        return self.metrics.copy()

# Utility functions
async def process_document_async(file_path: str, title: Optional[str] = None) -> ProcessedDocument:
    """Convenience function to process a document"""
    processor = DocumentProcessor()
    return await processor.process_document(file_path, title)

def estimate_tokens(text: str) -> int:
    """Estimate token count for text"""
    tokenizer = tiktoken.get_encoding("cl100k_base")
    return len(tokenizer.encode(text))

async def test_document_processor():
    """Test function for document processor"""
    processor = DocumentProcessor()
    
    # Test with sample text
    test_content = "This is a test document. " * 100
    
    with open("test_document.txt", "w") as f:
        f.write(test_content)
    
    try:
        result = await processor.process_document("test_document.txt", "Test Document")
        print(f"Processed document: {result.title}")
        print(f"Chunks: {len(result.chunks)}")
        print(f"Total tokens: {result.total_tokens}")
        print(f"Processing time: {result.processing_time_ms}ms")
    except Exception as e:
        print(f"Error: {e}")
    finally:
        # Cleanup
        Path("test_document.txt").unlink(missing_ok=True)

if __name__ == "__main__":
    asyncio.run(test_document_processor())
